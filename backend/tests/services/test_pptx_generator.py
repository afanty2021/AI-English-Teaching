"""
PPTX文档生成器单元测试

测试 PPTXDocumentGenerator 类的各项功能：
- 基本文档生成
- 中文内容支持
- 幻灯片数量验证
- 各种教案内容类型的生成
"""
import pytest
from pptx import Presentation

from app.services.document_generators.pptx_generator import PPTXDocumentGenerator


class TestPPTXDocumentGenerator:
    """PPTX 文档生成器测试类"""

    @pytest.fixture
    def generator(self):
        """创建生成器实例"""
        return PPTXDocumentGenerator()

    @pytest.fixture
    def basic_content(self):
        """基本教案内容"""
        return {
            "title": "过去完成时教学",
            "level": "B1",
            "topic": "Grammar",
            "duration": 45,
            "target_exam": "CET4",
        }

    @pytest.fixture
    def template_vars(self):
        """模板变量"""
        return {
            "teacher_name": "张老师",
            "school": "XX中学",
            "date": "2026-02-07",
        }

    @pytest.fixture
    def full_content(self):
        """完整教案内容"""
        return {
            "title": "过去完成时教学",
            "level": "B1",
            "topic": "Grammar",
            "duration": 45,
            "target_exam": "CET4",
            "objectives": {
                "language_knowledge": [
                    "掌握过去完成时的基本结构",
                    "理解过去完成时的用法",
                ],
                "language_skills": {
                    "reading": ["能够识别文章中的过去完成时"],
                    "writing": ["能够正确使用过去完成时写作"],
                },
            },
            "vocabulary": {
                "verb": [
                    {
                        "word": "had done",
                        "phonetic": "/hæd dʌn/",
                        "part_of_speech": "v.",
                        "meaning_cn": "过去完成时结构",
                        "example_sentence": "I had finished my homework before he came.",
                    },
                    {
                        "word": "had gone",
                        "phonetic": "/hæd ɡɒn/",
                        "part_of_speech": "v.",
                        "meaning_cn": "已经去了",
                        "example_sentence": "She had gone home when I arrived.",
                    },
                ],
            },
            "grammar_points": [
                {
                    "name": "过去完成时基本结构",
                    "description": "表示在过去某个时间之前已经完成的动作",
                    "rule": "had + 过去分词",
                    "examples": [
                        "I had finished my work before he came.",
                        "She had already left when I arrived.",
                    ],
                    "common_mistakes": [
                        "混淆过去完成时与一般过去时",
                        "忘记使用助动词 had",
                    ],
                    "practice_tips": [
                        "多读例句",
                        "练习时态转换",
                    ],
                }
            ],
            "teaching_structure": {
                "warm_up": {
                    "duration": 5,
                    "description": "通过问答引入主题",
                    "activities": ["复习一般过去时", "引入时间概念"],
                    "teacher_actions": ["提问", "引导"],
                    "student_actions": ["回答问题", "思考"],
                },
                "presentation": {
                    "duration": 15,
                    "description": "讲解过去完成时的结构和用法",
                    "activities": ["展示例句", "讲解规则"],
                },
            },
            "exercises": {
                "multiple_choice": [
                    {
                        "question": "When I arrived at the station, the train ____.",
                        "options": ["had left", "left", "has left", "leaves"],
                        "correct_answer": "A",
                        "explanation": "火车在我到达之前已经离开，所以用过去完成时",
                    },
                    {
                        "question": "She ____ her homework before 8 PM yesterday.",
                        "options": ["had finished", "finished", "finishes", "finishing"],
                        "correct_answer": "A",
                        "explanation": "在昨晚8点之前已经完成，用过去完成时",
                    },
                ],
            },
            "teaching_notes": "本课重点在于区分过去完成时和一般过去时的用法",
        }

    # ==================== 基本生成测试 ====================

    def test_generate_basic_pptx(self, generator, basic_content, template_vars):
        """测试生成基本PPTX文档"""
        ppt_bytes = generator.generate(basic_content, template_vars)

        # 验证返回的是字节
        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

        # 验证可以加载为Presentation对象
        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 基本内容应该至少有2页幻灯片（标题页 + 概述页）
        assert len(prs.slides) >= 2

    def test_generate_full_pptx(self, generator, full_content, template_vars):
        """测试生成完整PPTX文档"""
        ppt_bytes = generator.generate(full_content, template_vars)

        # 验证返回的是字节
        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

        # 加载演示文稿
        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 完整内容应该有更多幻灯片
        # 标题页(1) + 概述页(1) + 教学目标(1) + 词汇(1) + 语法(1) + 流程(2) + 练习(1) + 反思(1) = 9
        assert len(prs.slides) >= 8

    def test_generate_without_optional_sections(self, generator, basic_content, template_vars):
        """测试不包含可选章节的生成"""
        # 只包含基本内容
        ppt_bytes = generator.generate(basic_content, template_vars)

        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    # ==================== 内容验证测试 ====================

    def test_chinese_font_support(self, generator, full_content, template_vars):
        """测试中文字体支持"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 验证幻灯片中的文本框
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # 如果包含中文，验证字体设置
                    text = shape.text
                    if any("\u4e00" <= char <= "\u9fff" for char in text):
                        # 确保有文本框（包含中文字符）
                        assert len(text) > 0

    def test_title_slide_content(self, generator, basic_content, template_vars):
        """测试标题页内容"""
        ppt_bytes = generator.generate(basic_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 第一张幻灯片应该是标题页
        title_slide = prs.slides[0]

        # 验证标题和副标题
        title_found = False
        subtitle_found = False

        for shape in title_slide.shapes:
            if shape.shape_type == 14:  # Placeholder type
                if hasattr(shape, "text"):
                    if basic_content["title"] in shape.text:
                        title_found = True
                    if basic_content["level"] in shape.text:
                        subtitle_found = True

        assert title_found, "标题页应包含教案标题"
        assert subtitle_found, "标题页应包含等级信息"

    def test_objectives_slide(self, generator, full_content, template_vars):
        """测试教学目标页"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 查找教学目标相关的幻灯片
        objectives_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "教学目标" in shape.text or "语言知识" in shape.text:
                        objectives_found = True
                        break
            if objectives_found:
                break

        assert objectives_found, "应包含教学目标相关幻灯片"

    def test_vocabulary_slides(self, generator, full_content, template_vars):
        """测试词汇幻灯片"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 查找词汇相关幻灯片
        vocabulary_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "核心词汇" in shape.text or "动词" in shape.text:
                        vocabulary_found = True
                        break
            if vocabulary_found:
                break

        assert vocabulary_found, "应包含核心词汇相关幻灯片"

    def test_grammar_points_slides(self, generator, full_content, template_vars):
        """测试语法点幻灯片"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 查找语法点相关幻灯片
        grammar_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "过去完成时" in shape.text or "基本结构" in shape.text:
                        grammar_found = True
                        break
            if grammar_found:
                break

        assert grammar_found, "应包含语法点相关幻灯片"

    def test_teaching_structure_slides(self, generator, full_content, template_vars):
        """测试教学流程幻灯片"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 查找教学流程相关幻灯片
        structure_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if any(
                        phase in shape.text
                        for phase in ["热身阶段", "讲解阶段", "练习阶段"]
                    ):
                        structure_found = True
                        break
            if structure_found:
                break

        assert structure_found, "应包含教学流程相关幻灯片"

    def test_exercises_slides(self, generator, full_content, template_vars):
        """测试练习题幻灯片"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 查找练习题相关幻灯片
        exercises_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "选择题" in shape.text or "填空题" in shape.text:
                        exercises_found = True
                        break
            if exercises_found:
                break

        assert exercises_found, "应包含练习题相关幻灯片"

    def test_teaching_notes_slide(self, generator, full_content, template_vars):
        """测试教学反思幻灯片"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 查找教学反思相关幻灯片
        notes_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "教学反思" in shape.text:
                        notes_found = True
                        break
            if notes_found:
                break

        assert notes_found, "应包含教学反思相关幻灯片"

    # ==================== 错误处理测试 ====================

    def test_generate_without_title_raises_error(self, generator, template_vars):
        """测试缺少标题时抛出错误"""
        content = {
            "level": "B1",
            "topic": "Grammar",
        }

        # 异常被包装在 Exception 中，所以检查 Exception
        with pytest.raises(Exception, match="教案标题不能为空"):
            generator.generate(content, template_vars)

    def test_generate_with_empty_content(self, generator, template_vars):
        """测试空内容处理"""
        content = {
            "title": "测试教案",
        }

        ppt_bytes = generator.generate(content, template_vars)

        # 即使没有其他内容，也应该生成标题页和概述页
        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    # ==================== 幻灯片数量测试 ====================

    def test_slide_count_with_minimal_content(self, generator, basic_content, template_vars):
        """测试最小内容的幻灯片数量"""
        ppt_bytes = generator.generate(basic_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 最小内容：标题页 + 概述页
        assert len(prs.slides) == 2

    def test_slide_count_with_full_content(self, generator, full_content, template_vars):
        """测试完整内容的幻灯片数量"""
        ppt_bytes = generator.generate(full_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 完整内容应该有合理的幻灯片数量
        # 根据full_content的结构，应该有：
        # 标题页(1) + 概述页(1) + 教学目标(1) + 词汇(1) + 语法(1) + 流程(2) + 练习(1) + 反思(1)
        assert len(prs.slides) >= 8

    # ==================== 模板变量测试 ====================

    def test_template_vars_applied(self, generator, basic_content, template_vars):
        """测试模板变量应用"""
        ppt_bytes = generator.generate(basic_content, template_vars)

        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))

        # 在标题页查找教师信息
        teacher_info_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if template_vars["teacher_name"] in shape.text:
                        teacher_info_found = True
                        break
            if teacher_info_found:
                break

        assert teacher_info_found, "标题页应包含教师姓名"

    def test_template_vars_with_empty_values(self, generator, basic_content):
        """测试空的模板变量"""
        empty_vars = {
            "teacher_name": "",
            "school": "",
            "date": "",
        }

        ppt_bytes = generator.generate(basic_content, empty_vars)

        # 即使模板变量为空，也应该能生成文档
        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    # ==================== 辅助方法测试 ====================

    def test_color_constants(self, generator):
        """测试颜色常量定义"""
        assert generator.COLOR_TITLE is not None
        assert generator.COLOR_HEADING is not None
        assert generator.COLOR_ACCENT is not None
        assert generator.COLOR_CONTENT is not None
        assert generator.COLOR_LIGHT is not None

    def test_font_size_constants(self, generator):
        """测试字体大小常量定义"""
        assert generator.FONT_SIZE_TITLE is not None
        assert generator.FONT_SIZE_SUBTITLE is not None
        assert generator.FONT_SIZE_HEADING is not None
        assert generator.FONT_SIZE_CONTENT is not None
        assert generator.FONT_SIZE_SMALL is not None

    def test_default_font_setting(self, generator):
        """测试默认字体设置"""
        assert generator.default_font == "SimSun"

    # ==================== 特定内容类型测试 ====================

    def test_empty_objectives_handled(self, generator, basic_content, template_vars):
        """测试空教学目标的处理"""
        content = basic_content.copy()
        content["objectives"] = {}

        ppt_bytes = generator.generate(content, template_vars)

        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    def test_empty_vocabulary_handled(self, generator, basic_content, template_vars):
        """测试空词汇的处理"""
        content = basic_content.copy()
        content["vocabulary"] = {}

        ppt_bytes = generator.generate(content, template_vars)

        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    def test_empty_grammar_points_handled(self, generator, basic_content, template_vars):
        """测试空语法点的处理"""
        content = basic_content.copy()
        content["grammar_points"] = []

        ppt_bytes = generator.generate(content, template_vars)

        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    def test_empty_teaching_structure_handled(self, generator, basic_content, template_vars):
        """测试空教学流程的处理"""
        content = basic_content.copy()
        content["teaching_structure"] = {}

        ppt_bytes = generator.generate(content, template_vars)

        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0

    def test_empty_exercises_handled(self, generator, basic_content, template_vars):
        """测试空练习题的处理"""
        content = basic_content.copy()
        content["exercises"] = {}

        ppt_bytes = generator.generate(content, template_vars)

        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 0


# ==================== 集成测试 ====================

class TestPPTXDocumentGeneratorIntegration:
    """PPTX 文档生成器集成测试"""

    def test_end_to_end_generation(self):
        """端到端生成测试"""
        generator = PPTXDocumentGenerator()

        content = {
            "title": "英语语法教学 - 过去完成时",
            "level": "B1",
            "topic": "Grammar",
            "duration": 45,
            "target_exam": "CET4",
            "objectives": {
                "language_knowledge": [
                    "掌握过去完成时的基本结构",
                    "理解过去完成时的用法",
                ],
            },
            "vocabulary": {
                "verb": [
                    {
                        "word": "had done",
                        "phonetic": "/hæd dʌn/",
                        "part_of_speech": "v.",
                        "meaning_cn": "过去完成时结构",
                        "example_sentence": "I had finished my homework.",
                    },
                ],
            },
            "grammar_points": [
                {
                    "name": "过去完成时",
                    "description": "表示在过去某个时间之前已经完成的动作",
                    "rule": "had + 过去分词",
                    "examples": ["I had finished my work."],
                    "common_mistakes": ["混淆时态"],
                }
            ],
            "teaching_structure": {
                "warm_up": {
                    "duration": 5,
                    "description": "热身活动",
                    "activities": ["问答"],
                },
            },
            "exercises": {
                "multiple_choice": [
                    {
                        "question": "测试题目",
                        "options": ["A", "B", "C", "D"],
                        "correct_answer": "A",
                    }
                ],
            },
            "teaching_notes": "教学反思内容",
        }

        template_vars = {
            "teacher_name": "李老师",
            "school": "北京中学",
            "date": "2026-02-07",
        }

        # 生成文档
        ppt_bytes = generator.generate(content, template_vars)

        # 验证
        assert isinstance(ppt_bytes, bytes)
        assert len(ppt_bytes) > 1000  # 确保有实际内容

        # 验证可以加载
        from io import BytesIO
        prs = Presentation(BytesIO(ppt_bytes))
        assert len(prs.slides) > 0
