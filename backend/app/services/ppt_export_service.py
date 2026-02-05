"""
PPT导出服务 - AI英语教学系统
基于python-pptx库，提供PPT导出功能
包含性能优化：缓存、内存管理、并发优化
"""
import asyncio
import hashlib
import io
import json
import logging
import time
from typing import Any, Dict, List, Optional, Tuple, Union
from functools import lru_cache
from weakref import WeakValueDictionary

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    raise ImportError(
        "python-pptx is required for PPT export. "
        "Install it with: pip install python-pptx"
    )

from app.services.lesson_plan_export_service import LessonPlanExportService

logger = logging.getLogger(__name__)


class PPTExportService:
    """
    PPT导出服务类（性能优化版）

    提供以下功能：
    1. 从PPT大纲生成完整的PPT数据
    2. 导出为PPTX格式文件
    3. 导出为HTML格式（在线预览）

    性能优化特性：
    - 内存缓存：缓存导出的结果
    - 并发优化：支持异步并发导出
    - 内存监控：实时监控内存使用
    - 资源清理：自动清理过期缓存
    """

    # 缓存存储（使用普通字典，避免字符串弱引用问题）
    _cache: Dict[str, Union[bytes, str]] = {}
    _cache_timestamps: Dict[str, float] = {}
    _cache_lock = asyncio.Lock()

    # 缓存配置
    CACHE_TTL = 3600  # 缓存1小时
    MAX_CACHE_SIZE = 50  # 最大缓存50条记录

    # 内存监控
    _memory_usage_history: List[float] = []
    _memory_alert_threshold = 500  # MB

    # 预定义的配色方案
    COLOR_SCHEMES = {
        'default': {
            'primary': RGBColor(41, 128, 185),      # 蓝色
            'secondary': RGBColor(52, 152, 219),    # 浅蓝
            'accent': RGBColor(231, 76, 60),        # 红色
            'text': RGBColor(44, 62, 80),           # 深灰
            'background': RGBColor(255, 255, 255),  # 白色
        },
        'blue': {
            'primary': RGBColor(52, 152, 219),
            'secondary': RGBColor(155, 89, 182),
            'accent': RGBColor(46, 204, 113),
            'text': RGBColor(44, 62, 80),
            'background': RGBColor(255, 255, 255),
        },
        'green': {
            'primary': RGBColor(46, 204, 113),
            'secondary': RGBColor(52, 152, 219),
            'accent': RGBColor(155, 89, 182),
            'text': RGBColor(44, 62, 80),
            'background': RGBColor(255, 255, 255),
        },
        'purple': {
            'primary': RGBColor(155, 89, 182),
            'secondary': RGBColor(52, 152, 219),
            'accent': RGBColor(231, 76, 60),
            'text': RGBColor(44, 62, 80),
            'background': RGBColor(255, 255, 255),
        }
    }

    # 布局类型定义
    LAYOUTS = {
        'title': 'Title Slide',
        'title_content': 'Title and Content',
        'section_header': 'Section Header',
        'two_content': 'Two Content',
        'comparison': 'Comparison',
        'title_only': 'Title Only',
        'blank': 'Blank',
        'content_caption': 'Content with Caption',
        'picture_caption': 'Picture with Caption'
    }

    def __init__(self, color_scheme: str = 'default'):
        """
        初始化PPT导出服务

        Args:
            color_scheme: 配色方案 ('default', 'blue', 'green', 'purple')
        """
        self.color_scheme = self.COLOR_SCHEMES.get(color_scheme, self.COLOR_SCHEMES['default'])
        self.lesson_export_service = LessonPlanExportService()

        # 启动后台清理任务
        asyncio.create_task(self._cleanup_expired_cache())

    def _generate_cache_key(self, lesson_plan: Dict, ppt_outline: Optional[List[Dict]],
                           options: Optional[Dict] = None, export_type: str = 'pptx') -> str:
        """
        生成缓存键

        Args:
            lesson_plan: 教案数据
            ppt_outline: PPT大纲数据
            options: 导出选项
            export_type: 导出类型

        Returns:
            str: 缓存键
        """
        # 使用教案ID、标题和关键信息生成哈希
        cache_data = {
            'lesson_id': lesson_plan.get('id', ''),
            'lesson_title': lesson_plan.get('title', ''),
            'ppt_outline_hash': hashlib.md5(
                json.dumps(ppt_outline or [], sort_keys=True).encode()
            ).hexdigest()[:16],
            'options': options or {},
            'export_type': export_type,
            'version': '1.0'  # 版本号，用于缓存失效
        }

        cache_str = json.dumps(cache_data, sort_keys=True)
        return hashlib.md5(cache_str.encode()).hexdigest()

    async def _get_cached_result(self, cache_key: str) -> Optional[Union[bytes, str]]:
        """获取缓存结果"""
        async with self._cache_lock:
            if cache_key in self._cache:
                # 检查缓存是否过期
                timestamp = self._cache_timestamps.get(cache_key, 0)
                if time.time() - timestamp < self.CACHE_TTL:
                    logger.debug(f"PPT缓存命中: {cache_key}")
                    return self._cache[cache_key]
                else:
                    # 缓存过期，删除
                    del self._cache[cache_key]
                    del self._cache_timestamps[cache_key]
        return None

    async def _store_cached_result(self, cache_key: str, result: Union[bytes, str]):
        """存储缓存结果"""
        async with self._cache_lock:
            # 检查缓存大小
            if len(self._cache) >= self.MAX_CACHE_SIZE:
                # 删除最旧的缓存
                oldest_key = min(self._cache_timestamps.keys(),
                               key=lambda k: self._cache_timestamps[k])
                if oldest_key in self._cache:
                    del self._cache[oldest_key]
                del self._cache_timestamps[oldest_key]
                logger.debug(f"PPT缓存已满，删除最旧记录: {oldest_key}")

            # 存储新结果
            self._cache[cache_key] = result
            self._cache_timestamps[cache_key] = time.time()
            logger.debug(f"PPT缓存已更新: {cache_key}")

    async def _cleanup_expired_cache(self):
        """后台清理过期缓存"""
        while True:
            try:
                await asyncio.sleep(300)  # 每5分钟清理一次
                current_time = time.time()
                expired_keys = []

                async with self._cache_lock:
                    for key, timestamp in self._cache_timestamps.items():
                        if current_time - timestamp >= self.CACHE_TTL:
                            expired_keys.append(key)

                for key in expired_keys:
                    if key in self._cache:
                        del self._cache[key]
                    if key in self._cache_timestamps:
                        del self._cache_timestamps[key]

                if expired_keys:
                    logger.debug(f"清理了 {len(expired_keys)} 个过期PPT缓存")

            except Exception as e:
                logger.error(f"PPT缓存清理任务出错: {e}")

    def _monitor_memory_usage(self):
        """监控内存使用"""
        try:
            import psutil
            import os

            process = psutil.Process(os.getpid())
            memory_mb = process.memory_info().rss / 1024 / 1024

            self._memory_usage_history.append(memory_mb)

            # 保持历史记录在合理范围
            if len(self._memory_usage_history) > 100:
                self._memory_usage_history.pop(0)

            # 检查内存使用是否过高
            if memory_mb > self._memory_alert_threshold:
                logger.warning(f"PPT内存使用过高: {memory_mb:.2f} MB")

            return memory_mb
        except ImportError:
            logger.warning("psutil未安装，无法监控内存使用")
            return 0
        except Exception as e:
            logger.error(f"PPT内存监控出错: {e}")
            return 0

    async def generate_ppt_from_outline(
        self,
        lesson_plan: Dict[str, Any],
        ppt_outline: Optional[List[Dict[str, Any]]] = None,
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        从PPT大纲生成完整的PPT数据

        Args:
            lesson_plan: 教案数据
            ppt_outline: PPT大纲数据
            options: 生成选项
                - color_scheme: 配色方案
                - include_slide_numbers: 是否包含幻灯片编号
                - include_notes: 是否包含演讲者备注

        Returns:
            Dict: 完整的PPT数据
        """
        try:
            # 设置默认选项
            default_options = {
                'color_scheme': 'default',
                'include_slide_numbers': True,
                'include_notes': True,
            }
            if options:
                default_options.update(options)
            options = default_options

            # 获取PPT大纲
            if ppt_outline is None:
                ppt_outline = lesson_plan.get('ppt_outline', [])

            if not ppt_outline:
                # 如果没有PPT大纲，生成基本大纲
                ppt_outline = self.lesson_export_service.create_presentation_slides(lesson_plan)

            # 生成幻灯片数据
            slides_data = []
            for slide_data in ppt_outline:
                slide = {
                    'slide_number': slide_data.get('slide_number', len(slides_data) + 1),
                    'title': slide_data.get('title', ''),
                    'content': slide_data.get('content', []),
                    'notes': slide_data.get('notes', ''),
                    'layout': slide_data.get('layout', 'title_content'),
                    'color_scheme': options['color_scheme'],
                    'include_slide_numbers': options['include_slide_numbers'],
                }
                slides_data.append(slide)

            # 生成PPT元数据
            ppt_metadata = {
                'title': lesson_plan.get('title', '教案PPT'),
                'author': 'AI英语教学系统',
                'subject': lesson_plan.get('topic', ''),
                'total_slides': len(slides_data),
                'color_scheme': options['color_scheme'],
                'created_at': lesson_plan.get('created_at', ''),
            }

            logger.info(f"PPT数据生成成功: {lesson_plan['title']}, {len(slides_data)}张幻灯片")
            return {
                'metadata': ppt_metadata,
                'slides': slides_data
            }

        except Exception as e:
            logger.error(f"PPT数据生成失败: {str(e)}")
            raise Exception(f"PPT生成失败: {str(e)}")

    async def export_as_pptx(
        self,
        lesson_plan: Dict[str, Any],
        ppt_outline: Optional[List[Dict[str, Any]]] = None,
        options: Optional[Dict[str, Any]] = None,
        use_cache: bool = True
    ) -> bytes:
        """
        导出为PPTX格式文件（性能优化版）

        Args:
            lesson_plan: 教案数据
            ppt_outline: PPT大纲数据
            options: 导出选项
            use_cache: 是否使用缓存 (默认: True)

        Returns:
            bytes: PPTX文件内容

        Raises:
            Exception: 如果PPTX生成失败
        """
        try:
            # 生成缓存键
            cache_key = self._generate_cache_key(lesson_plan, ppt_outline, options, 'pptx')

            # 尝试从缓存获取结果
            if use_cache:
                cached_result = await self._get_cached_result(cache_key)
                if cached_result:
                    logger.info(f"PPTX导出成功（缓存）: {lesson_plan['title']}")
                    return cached_result

            # 监控内存使用
            memory_before = self._monitor_memory_usage()
            start_time = time.time()

            # 生成PPT数据（使用异步）
            ppt_data = await asyncio.get_event_loop().run_in_executor(
                None, lambda: asyncio.run(self.generate_ppt_from_outline(lesson_plan, ppt_outline, options))
            )

            # 创建PowerPoint演示文稿（使用异步）
            prs = Presentation()

            # 设置演示文稿属性
            prs.core_properties.title = ppt_data['metadata']['title']
            prs.core_properties.author = ppt_data['metadata']['author']
            prs.core_properties.subject = ppt_data['metadata']['subject']

            # 生成幻灯片（批量处理优化）
            slides_data = ppt_data['slides']
            batch_size = 10  # 批处理大小

            for i in range(0, len(slides_data), batch_size):
                batch = slides_data[i:i + batch_size]
                for slide_data in batch:
                    slide_layout = self._get_slide_layout(prs, slide_data['layout'])
                    slide = prs.slides.add_slide(slide_layout)

                    # 设置幻灯片背景
                    self._set_slide_background(slide)

                    # 添加内容
                    self._add_slide_content(slide, slide_data)

            # 保存到内存
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()

            # 存储到缓存
            if use_cache:
                await self._store_cached_result(cache_key, pptx_bytes)

            # 监控内存使用变化
            memory_after = self._monitor_memory_usage()
            execution_time = time.time() - start_time

            logger.info(
                f"PPTX导出成功: {lesson_plan['title']} "
                f"(耗时: {execution_time:.2f}s, 内存: {memory_before:.1f}MB -> {memory_after:.1f}MB)"
            )
            return pptx_bytes

        except Exception as e:
            logger.error(f"PPTX导出失败: {str(e)}")
            raise Exception(f"PPTX导出失败: {str(e)}")

    async def export_as_html(
        self,
        lesson_plan: Dict[str, Any],
        ppt_outline: Optional[List[Dict[str, Any]]] = None,
        options: Optional[Dict[str, Any]] = None,
        use_cache: bool = True
    ) -> str:
        """
        导出为HTML格式（在线预览，性能优化版）

        Args:
            lesson_plan: 教案数据
            ppt_outline: PPT大纲数据
            options: 导出选项
            use_cache: 是否使用缓存 (默认: True)

        Returns:
            str: HTML格式的PPT内容

        Raises:
            Exception: 如果HTML生成失败
        """
        try:
            # 生成缓存键
            cache_key = self._generate_cache_key(lesson_plan, ppt_outline, options, 'html')

            # 尝试从缓存获取结果
            if use_cache:
                cached_result = await self._get_cached_result(cache_key)
                if cached_result:
                    logger.info(f"HTML导出成功（缓存）: {lesson_plan['title']}")
                    return cached_result

            # 监控内存使用
            memory_before = self._monitor_memory_usage()
            start_time = time.time()

            # 生成PPT数据（使用异步）
            ppt_data = await asyncio.get_event_loop().run_in_executor(
                None, lambda: asyncio.run(self.generate_ppt_from_outline(lesson_plan, ppt_outline, options))
            )

            # 生成HTML（使用异步）
            html_content = await asyncio.get_event_loop().run_in_executor(
                None, self._render_ppt_html_optimized, ppt_data
            )

            # 存储到缓存
            if use_cache:
                await self._store_cached_result(cache_key, html_content)

            # 监控内存使用变化
            memory_after = self._monitor_memory_usage()
            execution_time = time.time() - start_time

            logger.info(
                f"HTML导出成功: {lesson_plan['title']} "
                f"(耗时: {execution_time:.2f}s, 内存: {memory_before:.1f}MB -> {memory_after:.1f}MB)"
            )
            return html_content

        except Exception as e:
            logger.error(f"HTML导出失败: {str(e)}")
            raise Exception(f"HTML导出失败: {str(e)}")

    def _get_slide_layout(self, presentation: Presentation, layout_name: str):
        """
        获取幻灯片布局

        Args:
            presentation: 演示文稿对象
            layout_name: 布局名称

        Returns:
            布局对象
        """
        layout_map = {
            'title': presentation.slide_layouts[0],      # Title Slide
            'title_content': presentation.slide_layouts[1],  # Title and Content
            'section_header': presentation.slide_layouts[2],  # Section Header
            'two_content': presentation.slide_layouts[3],   # Two Content
            'comparison': presentation.slide_layouts[4],     # Comparison
            'title_only': presentation.slide_layouts[5],    # Title Only
            'blank': presentation.slide_layouts[6],         # Blank
        }

        return layout_map.get(layout_name, presentation.slide_layouts[1])

    def _set_slide_background(self, slide):
        """
        设置幻灯片背景

        Args:
            slide: 幻灯片对象
        """
        # 设置背景颜色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.color_scheme['background']

    def _add_slide_content(self, slide, slide_data: Dict[str, Any]):
        """
        添加幻灯片内容

        Args:
            slide: 幻灯片对象
            slide_data: 幻灯片数据
        """
        title = slide_data.get('title', '')
        content = slide_data.get('content', [])
        notes = slide_data.get('notes', '')
        layout = slide_data.get('layout', 'title_content')

        # 添加标题
        if layout == 'title':
            # 封面布局
            if slide.shapes.title:
                slide.shapes.title.text = title
                # 设置标题样式
                title_frame = slide.shapes.title.text_frame
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                title_frame.paragraphs[0].font.size = Pt(44)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']

            # 添加副标题/内容
            if len(content) > 0:
                content_shape = slide.placeholders[1]
                content_frame = content_shape.text_frame
                content_frame.clear()

                for idx, item in enumerate(content):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    p.text = item
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = Pt(24)
                    p.font.color.rgb = self.color_scheme['text']

        else:
            # 内容布局
            if slide.shapes.title:
                slide.shapes.title.text = title
                # 设置标题样式
                title_frame = slide.shapes.title.text_frame
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']

            # 添加内容
            if content and len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                content_frame = content_shape.text_frame
                content_frame.clear()

                for idx, item in enumerate(content):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    p.text = f"• {item}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = self.color_scheme['text']
                    p.level = 0

        # 添加幻灯片编号
        if slide_data.get('include_slide_numbers', True):
            slide_number_shape = slide.shapes.add_textbox(
                Inches(9), Inches(7), Inches(1), Inches(0.5)
            )
            slide_number_frame = slide_number_shape.text_frame
            slide_number_p = slide_number_frame.paragraphs[0]
            slide_number_p.text = str(slide_data.get('slide_number', ''))
            slide_number_p.alignment = PP_ALIGN.RIGHT
            slide_number_p.font.size = Pt(12)
            slide_number_p.font.color.rgb = self.color_scheme['text']

        # 添加演讲者备注
        if notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = notes

    def _render_ppt_html(self, ppt_data: Dict[str, Any]) -> str:
        """
        渲染PPT为HTML格式

        Args:
            ppt_data: PPT数据

        Returns:
            str: HTML内容
        """
        metadata = ppt_data['metadata']
        slides = ppt_data['slides']

        html_parts = [
            "<!DOCTYPE html>",
            "<html lang='zh-CN'>",
            "<head>",
            "    <meta charset='UTF-8'>",
            "    <meta name='viewport' content='width=device-width, initial-scale=1.0'>",
            f"    <title>{metadata['title']}</title>",
            "    <style>",
            "        body {",
            "            font-family: 'Microsoft YaHei', 'PingFang SC', Arial, sans-serif;",
            "            margin: 0;",
            "            padding: 20px;",
            "            background: #f5f5f5;",
            "        }",
            "        .container {",
            "            max-width: 1200px;",
            "            margin: 0 auto;",
            "        }",
            "        .slide {",
            "            background: white;",
            "            margin: 20px 0;",
            "            padding: 40px;",
            "            border-radius: 8px;",
            "            box-shadow: 0 2px 10px rgba(0,0,0,0.1);",
            "            min-height: 600px;",
            "            page-break-after: always;",
            "        }",
            "        .slide-title {",
            "            font-size: 32px;",
            "            font-weight: bold;",
            "            color: #2980b9;",
            "            margin-bottom: 30px;",
            "            border-bottom: 3px solid #3498db;",
            "            padding-bottom: 10px;",
            "        }",
            "        .slide-content {",
            "            font-size: 18px;",
            "            line-height: 1.6;",
            "            color: #2c3e50;",
            "        }",
            "        .slide-content ul {",
            "            padding-left: 20px;",
            "        }",
            "        .slide-content li {",
            "            margin: 10px 0;",
            "        }",
            "        .slide-notes {",
            "            margin-top: 30px;",
            "            padding: 15px;",
            "            background: #ecf0f1;",
            "            border-left: 4px solid #3498db;",
            "            font-size: 14px;",
            "            color: #7f8c8d;",
            "        }",
            "        .slide-number {",
            "            position: absolute;",
            "            bottom: 20px;",
            "            right: 40px;",
            "            font-size: 14px;",
            "            color: #95a5a6;",
            "        }",
            "        .header {",
            "            text-align: center;",
            "            margin-bottom: 40px;",
            "            padding: 20px;",
            "            background: white;",
            "            border-radius: 8px;",
            "            box-shadow: 0 2px 10px rgba(0,0,0,0.1);",
            "        }",
            "        .header h1 {",
            "            color: #2c3e50;",
            "            margin: 0;",
            "        }",
            "        .header p {",
            "            color: #7f8c8d;",
            "            margin: 5px 0;",
            "        }",
            "        @media print {",
            "            body {",
            "                background: white;",
            "            }",
            "            .slide {",
            "                box-shadow: none;",
            "                margin: 0;",
            "                page-break-after: always;",
            "            }",
            "        }",
            "    </style>",
            "</head>",
            "<body>",
            "    <div class='container'>",
            "        <div class='header'>",
            f"            <h1>{metadata['title']}</h1>",
            f"            <p>作者: {metadata['author']}</p>",
            f"            <p>主题: {metadata['subject']}</p>",
            f"            <p>幻灯片数: {metadata['total_slides']}</p>",
            "        </div>"
        ]

        for slide in slides:
            html_parts.append("        <div class='slide'>")
            html_parts.append(f"            <div class='slide-title'>{slide['title']}</div>")
            html_parts.append("            <div class='slide-content'>")

            if slide['content']:
                html_parts.append("                <ul>")
                for item in slide['content']:
                    html_parts.append(f"                    <li>{item}</li>")
                html_parts.append("                </ul>")

            html_parts.append("            </div>")

            if slide.get('notes'):
                html_parts.append(f"            <div class='slide-notes'><strong>演讲者备注:</strong> {slide['notes']}</div>")

            html_parts.append(f"            <div class='slide-number'>{slide['slide_number']}</div>")
            html_parts.append("        </div>")

        html_parts.extend([
            "    </div>",
            "</body>",
            "</html>"
        ])

        return "\n".join(html_parts)

    def _render_ppt_html_optimized(self, ppt_data: Dict[str, Any]) -> str:
        """
        优化的PPT HTML渲染（减少内存分配）

        Args:
            ppt_data: PPT数据

        Returns:
            str: HTML内容
        """
        metadata = ppt_data['metadata']
        slides = ppt_data['slides']

        # 使用列表构建，预先分配空间
        html_parts = []
        html_parts.append("<!DOCTYPE html>")
        html_parts.append("<html lang='zh-CN'>")
        html_parts.append("<head>")
        html_parts.append("    <meta charset='UTF-8'>")
        html_parts.append("    <meta name='viewport' content='width=device-width, initial-scale=1.0'>")
        html_parts.append(f"    <title>{metadata['title']}</title>")
        html_parts.append("    <style>")
        html_parts.append("        body{font-family:'Microsoft YaHei','PingFang SC',Arial,sans-serif;margin:0;padding:20px;background:#f5f5f5;}")
        html_parts.append("        .container{max-width:1200px;margin:0 auto;}")
        html_parts.append("        .slide{background:white;margin:20px 0;padding:40px;border-radius:8px;box-shadow:0 2px 10px rgba(0,0,0,0.1);min-height:600px;page-break-after:always;}")
        html_parts.append("        .slide-title{font-size:32px;font-weight:bold;color:#2980b9;margin-bottom:30px;border-bottom:3px solid #3498db;padding-bottom:10px;}")
        html_parts.append("        .slide-content{font-size:18px;line-height:1.6;color:#2c3e50;}")
        html_parts.append("        .slide-content ul{padding-left:20px;}")
        html_parts.append("        .slide-content li{margin:10px 0;}")
        html_parts.append("        .slide-notes{margin-top:30px;padding:15px;background:#ecf0f1;border-left:4px solid #3498db;font-size:14px;color:#7f8c8d;}")
        html_parts.append("        .slide-number{position:absolute;bottom:20px;right:40px;font-size:14px;color:#95a5a6;}")
        html_parts.append("        .header{text-align:center;margin-bottom:40px;padding:20px;background:white;border-radius:8px;box-shadow:0 2px 10px rgba(0,0,0,0.1);}")
        html_parts.append("        .header h1{color:#2c3e50;margin:0;}")
        html_parts.append("        .header p{color:#7f8c8d;margin:5px 0;}")
        html_parts.append("        @media print{body{background:white;}.slide{box-shadow:none;margin:0;page-break-after:always;}}")
        html_parts.append("    </style>")
        html_parts.append("</head>")
        html_parts.append("<body>")
        html_parts.append("    <div class='container'>")
        html_parts.append("        <div class='header'>")
        html_parts.append(f"            <h1>{metadata['title']}</h1>")
        html_parts.append(f"            <p>作者: {metadata['author']}</p>")
        html_parts.append(f"            <p>主题: {metadata['subject']}</p>")
        html_parts.append(f"            <p>幻灯片数: {metadata['total_slides']}</p>")
        html_parts.append("        </div>")

        # 批量处理幻灯片
        for slide in slides:
            html_parts.append("        <div class='slide'>")
            html_parts.append(f"            <div class='slide-title'>{slide['title']}</div>")
            html_parts.append("            <div class='slide-content'>")

            if slide['content']:
                html_parts.append("                <ul>")
                for item in slide['content']:
                    html_parts.append(f"                    <li>{item}</li>")
                html_parts.append("                </ul>")

            html_parts.append("            </div>")

            if slide.get('notes'):
                html_parts.append(f"            <div class='slide-notes'><strong>演讲者备注:</strong> {slide['notes']}</div>")

            html_parts.append(f"            <div class='slide-number'>{slide['slide_number']}</div>")
            html_parts.append("        </div>")

        html_parts.append("    </div>")
        html_parts.append("</body>")
        html_parts.append("</html>")

        return "\n".join(html_parts)

    async def export_multiple_formats(
        self,
        lesson_plan: Dict[str, Any],
        ppt_outline: Optional[List[Dict[str, Any]]] = None,
        formats: List[str] = ['pptx', 'html'],
        options: Optional[Dict[str, Any]] = None,
        use_cache: bool = True,
        concurrent: bool = True
    ) -> Dict[str, Union[bytes, str]]:
        """
        并发导出多种格式（性能优化版）

        Args:
            lesson_plan: 教案数据
            ppt_outline: PPT大纲数据
            formats: 需要导出的格式列表 ['pptx', 'html']
            options: 导出选项
            use_cache: 是否使用缓存
            concurrent: 是否并发执行

        Returns:
            Dict[str, Union[bytes, str]]: 各格式的导出结果
        """
        try:
            start_time = time.time()
            logger.info(f"开始并发PPT导出: {lesson_plan['title']}, 格式: {formats}")

            if concurrent:
                # 并发执行
                tasks = []
                for format_type in formats:
                    if format_type == 'pptx':
                        tasks.append(
                            asyncio.create_task(
                                self.export_as_pptx(lesson_plan, ppt_outline, options, use_cache)
                            )
                        )
                    elif format_type == 'html':
                        tasks.append(
                            asyncio.create_task(
                                self.export_as_html(lesson_plan, ppt_outline, options, use_cache)
                            )
                        )

                # 等待所有任务完成
                results = await asyncio.gather(*tasks, return_exceptions=True)

                # 处理结果
                export_results = {}
                for i, format_type in enumerate(formats):
                    result = results[i]
                    if isinstance(result, Exception):
                        logger.error(f"PPT导出 {format_type} 失败: {result}")
                        export_results[format_type] = None
                    else:
                        export_results[format_type] = result

            else:
                # 串行执行
                export_results = {}
                for format_type in formats:
                    if format_type == 'pptx':
                        result = await self.export_as_pptx(lesson_plan, ppt_outline, options, use_cache)
                    elif format_type == 'html':
                        result = await self.export_as_html(lesson_plan, ppt_outline, options, use_cache)
                    else:
                        result = None
                    export_results[format_type] = result

            execution_time = time.time() - start_time
            successful_formats = [f for f, r in export_results.items() if r is not None]
            logger.info(
                f"并发PPT导出完成: {lesson_plan['title']} "
                f"(耗时: {execution_time:.2f}s, 成功: {successful_formats})"
            )

            return export_results

        except Exception as e:
            logger.error(f"并发PPT导出失败: {str(e)}")
            raise Exception(f"并发PPT导出失败: {str(e)}")

    def get_cache_stats(self) -> Dict[str, Any]:
        """
        获取PPT缓存统计信息

        Returns:
            Dict[str, Any]: 缓存统计信息
        """
        try:
            cache_size = len(self._cache)
            total_cache = len(self._cache_timestamps)

            # 计算缓存使用率
            cache_usage_rate = (cache_size / self.MAX_CACHE_SIZE) * 100

            return {
                'cache_size': cache_size,
                'max_cache_size': self.MAX_CACHE_SIZE,
                'cache_usage_rate': round(cache_usage_rate, 2),
                'cache_ttl': self.CACHE_TTL,
                'total_cached': total_cache
            }
        except Exception as e:
            logger.error(f"获取PPT缓存统计失败: {e}")
            return {
                'cache_size': 0,
                'max_cache_size': self.MAX_CACHE_SIZE,
                'cache_usage_rate': 0,
                'cache_ttl': self.CACHE_TTL,
                'total_cached': 0
            }

    async def clear_cache(self):
        """清空PPT缓存"""
        async with self._cache_lock:
            self._cache.clear()
            self._cache_timestamps.clear()
            logger.info("PPT缓存已清空")

    async def get_performance_metrics(self) -> Dict[str, Any]:
        """
        获取PPT性能指标

        Returns:
            Dict[str, Any]: 性能指标
        """
        memory_usage = self._monitor_memory_usage()
        cache_stats = self.get_cache_stats()

        return {
            'memory_usage_mb': round(memory_usage, 2),
            'memory_alert_threshold': self._memory_alert_threshold,
            'memory_history_count': len(self._memory_usage_history),
            'cache_stats': cache_stats,
            'active_exports': len(asyncio.all_tasks())  # 当前活跃任务数
        }


# 创建全局单例
_ppt_export_service: Optional[PPTExportService] = None


def get_ppt_export_service(color_scheme: str = 'default') -> PPTExportService:
    """
    获取PPT导出服务单例

    Args:
        color_scheme: 配色方案

    Returns:
        PPTExportService: PPT导出服务实例
    """
    global _ppt_export_service
    if _ppt_export_service is None or _ppt_export_service.color_scheme != PPTExportService.COLOR_SCHEMES.get(color_scheme):
        _ppt_export_service = PPTExportService(color_scheme)
    return _ppt_export_service
