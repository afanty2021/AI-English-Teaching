"""
PPTX文档生成器 - AI英语教学系统

使用 python-pptx 库将教案数据转换为 PowerPoint 演示文稿。
支持中文内容、多种幻灯片布局、样式设置等功能。
"""
import logging
from io import BytesIO
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


class PPTXDocumentGenerator:
    """
    PPTX 文档生成器

    将教案数据转换为 PowerPoint (.pptx) 演示文稿。

    功能：
    - 生成包含教案完整内容的 PowerPoint 演示文稿
    - 支持中文内容（自动设置中文字体）
    - 使用多种幻灯片布局展示不同内容
    - 支持自定义样式和格式

    使用示例：
        ```python
        generator = PPTXDocumentGenerator()

        content = {
            "title": "过去完成时教学",
            "level": "B1",
            "topic": "Grammar",
            "duration": 45,
            # ... 其他教案内容
        }

        template_vars = {
            "teacher_name": "张老师",
            "school": "XX中学",
            "date": "2026-02-07"
        }

        ppt_bytes = generator.generate(content, template_vars)
        ```
    """

    # 样式配置
    # 字体大小（点）
    FONT_SIZE_TITLE = Pt(32)
    FONT_SIZE_SUBTITLE = Pt(20)
    FONT_SIZE_HEADING = Pt(24)
    FONT_SIZE_CONTENT = Pt(14)
    FONT_SIZE_SMALL = Pt(12)

    # 颜色配置
    COLOR_TITLE = RGBColor(0, 51, 102)      # 深蓝色
    COLOR_HEADING = RGBColor(0, 76, 153)    # 中蓝色
    COLOR_ACCENT = RGBColor(204, 0, 0)      # 红色用于强调
    COLOR_CONTENT = RGBColor(51, 51, 51)    # 深灰色
    COLOR_LIGHT = RGBColor(128, 128, 128)   # 浅灰色

    # 间距配置
    MARGIN_INCHES = Inches(0.5)
    SPACING_INCHES = Inches(0.15)

    def __init__(self):
        """初始化 PPTX 文档生成器"""
        self.prs: Optional[Presentation] = None
        self.default_font = "SimSun"  # 宋体，支持中文

    def generate(self, content: Dict[str, Any], template_vars: Dict[str, Any]) -> bytes:
        """
        生成完整的 PowerPoint 演示文稿

        Args:
            content: 教案内容数据，包含以下字段：
                - title: 教案标题
                - level: CEFR等级
                - topic: 主题
                - duration: 课程时长（分钟）
                - target_exam: 目标考试（可选）
                - objectives: 教学目标（字典）
                - vocabulary: 词汇表（字典）
                - grammar_points: 语法点（列表）
                - teaching_structure: 教学流程（字典）
                - leveled_materials: 分层材料（列表）
                - exercises: 练习题（字典）
                - ppt_outline: PPT大纲（列表）
                - teaching_notes: 教学反思（可选）
            template_vars: 模板变量，包含：
                - teacher_name: 教师姓名
                - school: 学校名称（可选）
                - date: 日期（可选）

        Returns:
            bytes: PowerPoint 演示文稿的二进制内容

        Raises:
            ValueError: 如果内容数据无效
            Exception: 如果文档生成失败
        """
        try:
            # 验证必要字段
            if not content.get("title"):
                raise ValueError("教案标题不能为空")

            # 创建新演示文稿
            self.prs = Presentation()

            # 添加标题页
            self._add_title_slide(content, template_vars)

            # 添加概述页
            self._add_overview_slide(content, template_vars)

            # 添加教学目标页
            if content.get("objectives"):
                self._add_objectives_slide(content["objectives"])

            # 添加核心词汇页
            if content.get("vocabulary"):
                self._add_vocabulary_slides(content["vocabulary"])

            # 添加语法点页
            if content.get("grammar_points"):
                self._add_grammar_points_slides(content["grammar_points"])

            # 添加教学流程页
            if content.get("teaching_structure"):
                self._add_teaching_structure_slides(content["teaching_structure"])

            # 添加分层材料页
            if content.get("leveled_materials"):
                self._add_leveled_materials_slides(content["leveled_materials"])

            # 添加练习题页
            if content.get("exercises"):
                self._add_exercises_slides(content["exercises"])

            # 添加教学反思页（如果有）
            if content.get("teaching_notes"):
                self._add_notes_slide(content["teaching_notes"])

            # 保存到字节流
            ppt_bytes = self._save_to_bytes()

            logger.info(f"PPTX文档生成成功: {content.get('title')}, 共 {len(self.prs.slides)} 页")
            return ppt_bytes

        except Exception as e:
            logger.error(f"PPTX文档生成失败: {str(e)}")
            raise Exception(f"PPTX文档生成失败: {str(e)}")

    def _save_to_bytes(self) -> bytes:
        """
        将演示文稿保存到字节流

        Returns:
            bytes: 演示文稿的二进制内容
        """
        ppt_stream = BytesIO()
        self.prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream.read()

    def _add_title_slide(self, content: Dict[str, Any], template_vars: Dict[str, Any]) -> None:
        """
        添加标题页

        Args:
            content: 教案内容
            template_vars: 模板变量
        """
        # 使用标题布局
        slide_layout = self.prs.slide_layouts[0]  # 标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = content.get("title", "教案")

        # 设置副标题
        subtitle = slide.placeholders[1]
        subtitle_text = f"{content.get('level', '')} 等级 - {content.get('topic', '')}"
        subtitle.text = subtitle_text

        # 添加教师信息（使用文本框）
        left = Inches(1)
        top = Inches(5)
        width = Inches(8)
        height = Inches(1)

        info_box = slide.shapes.add_textbox(left, top, width, height)
        info_frame = info_box.text_frame
        info_frame.word_wrap = True

        # 添加教师信息
        p = info_frame.paragraphs[0]
        p.text = f"授课教师: {template_vars.get('teacher_name', '')}"
        self._set_paragraph_style(p, font_size=self.FONT_SIZE_SMALL, color=self.COLOR_LIGHT)

        # 添加学校信息（如果有）
        if template_vars.get("school"):
            p = info_frame.add_paragraph()
            p.text = f"学校: {template_vars['school']}"
            self._set_paragraph_style(p, font_size=self.FONT_SIZE_SMALL, color=self.COLOR_LIGHT)

        # 添加日期信息（如果有）
        if template_vars.get("date"):
            p = info_frame.add_paragraph()
            p.text = f"日期: {template_vars['date']}"
            self._set_paragraph_style(p, font_size=self.FONT_SIZE_SMALL, color=self.COLOR_LIGHT)

    def _add_overview_slide(self, content: Dict[str, Any], template_vars: Dict[str, Any]) -> None:
        """
        添加概述页

        Args:
            content: 教案内容
            template_vars: 模板变量
        """
        slide = self._add_blank_slide()

        # 添加标题
        self._add_slide_title(slide, "课程概述")

        # 创建概述信息
        overview_items = [
            f"教案标题: {content.get('title', '')}",
            f"课程等级: {content.get('level', '')}",
            f"课程主题: {content.get('topic', '')}",
            f"课程时长: {content.get('duration', 0)} 分钟",
        ]

        if content.get("target_exam"):
            overview_items.append(f"目标考试: {content['target_exam']}")

        # 添加列表
        self._add_bulleted_list(slide, overview_items, top=Inches(2))

    def _add_objectives_slide(self, objectives: Dict[str, Any]) -> None:
        """
        添加教学目标页

        Args:
            objectives: 教学目标数据
        """
        # 类别名称映射
        category_names = {
            "language_knowledge": "语言知识",
            "language_skills": "语言技能",
            "learning_strategies": "学习策略",
            "cultural_awareness": "文化意识",
            "emotional_attitudes": "情感态度",
        }

        # 技能名称映射
        skill_names = {
            "listening": "听力",
            "speaking": "口语",
            "reading": "阅读",
            "writing": "写作",
        }

        for category, items in objectives.items():
            if not items:
                continue

            slide = self._add_blank_slide()

            # 添加类别标题
            category_name = category_names.get(category, category)
            self._add_slide_title(slide, f"教学目标 - {category_name}")

            # 添加内容
            y_position = Inches(2)

            if isinstance(items, dict):
                # 处理嵌套结构（如 language_skills）
                for skill, skill_items in items.items():
                    if skill_items:
                        skill_name = skill_names.get(skill, skill)
                        y_position = self._add_text_to_slide(
                            slide,
                            f"{skill_name}:",
                            top=y_position,
                            font_size=self.FONT_SIZE_CONTENT,
                            bold=True
                        )

                        for item in skill_items:
                            y_position = self._add_text_to_slide(
                                slide,
                                f"  • {item}",
                                top=y_position,
                                font_size=self.FONT_SIZE_CONTENT,
                                color=self.COLOR_CONTENT
                            )

            elif isinstance(items, list):
                # 处理列表结构
                for item in items:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"• {item}",
                        top=y_position,
                        font_size=self.FONT_SIZE_CONTENT,
                        color=self.COLOR_CONTENT
                    )

    def _add_vocabulary_slides(self, vocabulary: Dict[str, Any]) -> None:
        """
        添加核心词汇幻灯片

        Args:
            vocabulary: 词汇数据，按词性分组
        """
        # 词性名称映射
        pos_names = {
            "noun": "名词",
            "verb": "动词",
            "adj": "形容词",
            "adv": "副词",
            "prep": "介词",
            "phrase": "短语",
        }

        for pos_type, words in vocabulary.items():
            if not words:
                continue

            # 每种词性创建一个幻灯片
            slide = self._add_blank_slide()

            pos_name = pos_names.get(pos_type, pos_type)
            self._add_slide_title(slide, f"核心词汇 - {pos_name}")

            # 添加词汇列表
            y_position = Inches(2)

            for word_data in words:
                word = word_data.get("word", "")
                meaning = word_data.get("meaning_cn", "")
                phonetic = word_data.get("phonetic", "")
                example = word_data.get("example_sentence", "")

                # 单词和释义
                text = f"• {word}"
                if phonetic:
                    text += f" {phonetic}"
                text += f" - {meaning}"

                y_position = self._add_text_to_slide(
                    slide,
                    text,
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT,
                    bold=True
                )

                # 例句
                if example:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"  例: {example}",
                        top=y_position,
                        font_size=self.FONT_SIZE_SMALL,
                        color=self.COLOR_LIGHT
                    )

                # 间距
                y_position += Inches(0.1)

    def _add_grammar_points_slides(self, grammar_points: List[Dict[str, Any]]) -> None:
        """
        添加语法点幻灯片

        Args:
            grammar_points: 语法点列表
        """
        for gp in grammar_points:
            slide = self._add_blank_slide()

            # 语法点名称作为标题
            name = gp.get("name", "语法点")
            self._add_slide_title(slide, name)

            y_position = Inches(2)

            # 描述
            description = gp.get("description", "")
            if description:
                y_position = self._add_text_to_slide(
                    slide,
                    f"描述: {description}",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT
                )

            # 规则
            rule = gp.get("rule", "")
            if rule:
                y_position = self._add_text_to_slide(
                    slide,
                    f"规则: {rule}",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT,
                    color=self.COLOR_HEADING
                )

            # 例句
            examples = gp.get("examples", [])
            if examples:
                y_position = self._add_text_to_slide(
                    slide,
                    "例句:",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT,
                    bold=True
                )

                for example in examples:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"  • {example}",
                        top=y_position,
                        font_size=self.FONT_SIZE_CONTENT,
                        color=self.COLOR_CONTENT
                    )

            # 常见错误
            mistakes = gp.get("common_mistakes", [])
            if mistakes:
                y_position = self._add_text_to_slide(
                    slide,
                    "常见错误:",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT,
                    bold=True
                )

                for mistake in mistakes:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"  ✗ {mistake}",
                        top=y_position,
                        font_size=self.FONT_SIZE_SMALL,
                        color=self.COLOR_ACCENT
                    )

    def _add_teaching_structure_slides(self, structure: Dict[str, Any]) -> None:
        """
        添加教学流程幻灯片

        Args:
            structure: 教学流程数据
        """
        # 阶段名称映射
        phase_names = {
            "warm_up": "热身阶段",
            "presentation": "讲解阶段",
            "practice": "练习阶段",
            "production": "产出阶段",
            "summary": "总结阶段",
            "homework": "课后作业",
        }

        for phase, data in structure.items():
            if not data:
                continue

            slide = self._add_blank_slide()

            # 阶段标题
            phase_name = phase_names.get(phase, phase)
            self._add_slide_title(slide, phase_name)

            y_position = Inches(2)

            if isinstance(data, dict):
                # 时长
                duration = data.get("duration", 0)
                if duration:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"时长: {duration} 分钟",
                        top=y_position,
                        font_size=self.FONT_SIZE_CONTENT,
                        bold=True
                    )

                # 描述
                description = data.get("description", "")
                if description:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"描述: {description}",
                        top=y_position,
                        font_size=self.FONT_SIZE_CONTENT
                    )

                # 活动
                activities = data.get("activities", [])
                if activities:
                    y_position = self._add_text_to_slide(
                        slide,
                        "活动:",
                        top=y_position,
                        font_size=self.FONT_SIZE_CONTENT,
                        bold=True
                    )

                    for activity in activities:
                        y_position = self._add_text_to_slide(
                            slide,
                            f"  • {activity}",
                            top=y_position,
                            font_size=self.FONT_SIZE_CONTENT
                        )

                # 教师活动
                teacher_actions = data.get("teacher_actions", [])
                if teacher_actions:
                    y_position = self._add_text_to_slide(
                        slide,
                        "教师活动:",
                        top=y_position,
                        font_size=self.FONT_SIZE_CONTENT,
                        bold=True
                    )

                    for action in teacher_actions:
                        y_position = self._add_text_to_slide(
                            slide,
                            f"  • {action}",
                            top=y_position,
                            font_size=self.FONT_SIZE_CONTENT
                        )

    def _add_leveled_materials_slides(self, materials: List[Dict[str, Any]]) -> None:
        """
        添加分层材料幻灯片

        Args:
            materials: 分层材料列表
        """
        for material in materials:
            slide = self._add_blank_slide()

            # 材料标题
            title = material.get("title", "")
            level = material.get("level", "")
            slide_title = f"{title} (CEFR {level})" if title else f"CEFR {level}"
            self._add_slide_title(slide, slide_title)

            y_position = Inches(2)

            # 基本信息
            y_position = self._add_text_to_slide(
                slide,
                f"等级: {level} | 字数: {material.get('word_count', 0)} 字",
                top=y_position,
                font_size=self.FONT_SIZE_SMALL,
                color=self.COLOR_LIGHT
            )

            # 内容（截取前300字符）
            content = material.get("content", "")
            if content:
                preview = content[:300] + "..." if len(content) > 300 else content
                y_position = self._add_text_to_slide(
                    slide,
                    f"内容预览: {preview}",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT
                )

            # 重点词汇
            vocab_list = material.get("vocabulary_list", [])
            if vocab_list:
                y_position = self._add_text_to_slide(
                    slide,
                    "重点词汇:",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT,
                    bold=True
                )

                for vocab in vocab_list[:5]:  # 最多显示5个
                    y_position = self._add_text_to_slide(
                        slide,
                        f"  • {vocab.get('word', '')}: {vocab.get('meaning_cn', '')}",
                        top=y_position,
                        font_size=self.FONT_SIZE_SMALL
                    )

    def _add_exercises_slides(self, exercises: Dict[str, Any]) -> None:
        """
        添加练习题幻灯片

        Args:
            exercises: 练习题数据，按类型分组
        """
        # 题型名称映射
        type_names = {
            "multiple_choice": "选择题",
            "fill_blank": "填空题",
            "matching": "匹配题",
            "essay": "写作题",
            "speaking": "口语题",
            "translation": "翻译题",
        }

        for ex_type, exercise_list in exercises.items():
            if not exercise_list:
                continue

            # 每种题型创建一个幻灯片
            slide = self._add_blank_slide()

            type_name = type_names.get(ex_type, ex_type)
            self._add_slide_title(slide, type_name)

            # 添加题目（最多3题）
            y_position = Inches(2)

            for idx, exercise in enumerate(exercise_list[:3], 1):
                question = exercise.get("question", "")

                y_position = self._add_text_to_slide(
                    slide,
                    f"{idx}. {question}",
                    top=y_position,
                    font_size=self.FONT_SIZE_CONTENT,
                    bold=True
                )

                # 选项（选择题）
                options = exercise.get("options", [])
                if options:
                    option_labels = ["A", "B", "C", "D", "E", "F"]
                    for i, option in enumerate(options):
                        if i < len(option_labels):
                            y_position = self._add_text_to_slide(
                                slide,
                                f"  {option_labels[i]}. {option}",
                                top=y_position,
                                font_size=self.FONT_SIZE_SMALL
                            )

                # 答案
                answer = exercise.get("correct_answer", "")
                if answer:
                    y_position = self._add_text_to_slide(
                        slide,
                        f"  答案: {answer}",
                        top=y_position,
                        font_size=self.FONT_SIZE_SMALL,
                        color=self.COLOR_ACCENT
                    )

                # 间距
                y_position += Inches(0.1)

    def _add_notes_slide(self, notes: str) -> None:
        """
        添加教学反思页

        Args:
            notes: 教学反思内容
        """
        slide = self._add_blank_slide()

        self._add_slide_title(slide, "教学反思")

        # 添加反思内容
        self._add_text_to_slide(
            slide,
            notes,
            top=Inches(2),
            font_size=self.FONT_SIZE_CONTENT
        )

    # ==================== 辅助方法 ====================

    def _add_blank_slide(self):
        """添加空白幻灯片"""
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        return self.prs.slides.add_slide(blank_layout)

    def _add_slide_title(self, slide, title_text: str) -> None:
        """
        添加幻灯片标题

        Args:
            slide: 幻灯片对象
            title_text: 标题文本
        """
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.text = title_text
        self._set_paragraph_style(
            p,
            font_size=self.FONT_SIZE_HEADING,
            bold=True,
            color=self.COLOR_TITLE
        )

    def _add_text_to_slide(
        self,
        slide,
        text: str,
        top: Inches,
        font_size: Pt = FONT_SIZE_CONTENT,
        bold: bool = False,
        color: Optional[RGBColor] = None
    ) -> Inches:
        """
        在幻灯片上添加文本

        Args:
            slide: 幻灯片对象
            text: 文本内容
            top: 顶部位置
            font_size: 字体大小
            bold: 是否加粗
            color: 文字颜色

        Returns:
            Inches: 下一个文本的位置
        """
        left = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = text
        self._set_paragraph_style(
            p,
            font_size=font_size,
            bold=bold,
            color=color or self.COLOR_CONTENT
        )

        return top + height + Inches(0.05)

    def _add_bulleted_list(self, slide, items: List[str], top: Inches) -> None:
        """
        在幻灯片上添加项目符号列表

        Args:
            slide: 幻灯片对象
            items: 列表项
            top: 顶部位置
        """
        left = Inches(0.5)
        width = Inches(9)
        height = Inches(0.4)

        for i, item in enumerate(items):
            text_box = slide.shapes.add_textbox(
                left,
                top + Inches(i * 0.4),
                width,
                height
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.text = f"• {item}"
            p.level = 0
            self._set_paragraph_style(
                p,
                font_size=self.FONT_SIZE_CONTENT,
                color=self.COLOR_CONTENT
            )

    def _set_paragraph_style(
        self,
        paragraph,
        font_size: Pt = FONT_SIZE_CONTENT,
        bold: bool = False,
        color: RGBColor = COLOR_CONTENT
    ) -> None:
        """
        设置段落样式

        Args:
            paragraph: 段落对象
            font_size: 字体大小
            bold: 是否加粗
            color: 文字颜色
        """
        # 设置字体
        for run in paragraph.runs:
            run.font.name = self.default_font
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color
