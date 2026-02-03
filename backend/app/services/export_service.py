"""
教案导出服务 - AI英语教学系统

支持将教案导出为以下格式：
- Word (.docx)
- PowerPoint (.pptx)
- PDF (.pdf)
"""
import os
import io
import tempfile
from datetime import datetime
from typing import Optional
from uuid import UUID

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from app.models.lesson_plan import LessonPlan


class ExportService:
    """
    教案导出服务类

    支持将教案导出为 Word、PowerPoint 和 PDF 格式。
    """

    # 临时文件目录
    TEMP_DIR = tempfile.gettempdir()

    @staticmethod
    async def export_to_word(lesson_plan: LessonPlan) -> tuple[bytes, str]:
        """
        导出教案为 Word 文档

        Args:
            lesson_plan: 教案对象

        Returns:
            tuple[bytes, str]: (文件内容, 文件名)
        """
        doc = Document()

        # 设置文档标题
        title = doc.add_heading(lesson_plan.title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 添加基本信息
        info_table = doc.add_table(rows=5, cols=2)
        info_table.style = 'Light Grid Accent 1'

        info_data = [
            ("教学主题", lesson_plan.topic),
            ("CEFR 等级", lesson_plan.level),
            ("课程时长", f"{lesson_plan.duration} 分钟"),
            ("目标考试", lesson_plan.target_exam or "通用"),
            ("生成时间", lesson_plan.last_generated_at.strftime("%Y-%m-%d %H:%M") if lesson_plan.last_generated_at else ""),
        ]

        for i, (key, value) in enumerate(info_data):
            info_table.rows[i].cells[0].text = key
            info_table.rows[i].cells[1].text = str(value)

        doc.add_paragraph()

        # 教学目标
        if lesson_plan.objectives:
            doc.add_heading("一、教学目标", level=1)
            objectives = lesson_plan.objectives

            if objectives.get("language_knowledge"):
                doc.add_heading("知识目标", level=2)
                for obj in objectives["language_knowledge"]:
                    doc.add_paragraph(obj, style='List Bullet')

            if objectives.get("language_skills"):
                doc.add_heading("能力目标", level=2)
                skills = objectives["language_skills"]
                for skill_type, skill_list in skills.items():
                    if skill_list:
                        skill_name = {
                            "listening": "听力",
                            "speaking": "口语",
                            "reading": "阅读",
                            "writing": "写作"
                        }.get(skill_type, skill_type)
                        doc.add_paragraph(f"{skill_name}:", style='List Bullet')
                        for skill in skill_list:
                            doc.add_paragraph(f"  {skill}", style='List Bullet 2')

        # 核心词汇
        if lesson_plan.vocabulary:
            doc.add_heading("二、核心词汇", level=1)
            vocabulary = lesson_plan.vocabulary

            for pos_type, words in vocabulary.items():
                if words:
                    pos_name = {
                        "noun": "名词",
                        "verb": "动词",
                        "adjective": "形容词",
                        "adverb": "副词",
                        "phrase": "短语"
                    }.get(pos_type, pos_type)
                    doc.add_heading(pos_name, level=2)

                    for word in words[:20]:  # 限制数量
                        p = doc.add_paragraph()
                        word_text = word.get("word", "")
                        phonetic = word.get("phonetic", "")
                        meaning = word.get("meaning_cn", "")

                        # 添加单词
                        run = p.add_run(f"{word_text} ")
                        run.bold = True
                        run.font.size = Pt(12)

                        # 添加音标
                        if phonetic:
                            run = p.add_run(f"{phonetic} ")
                            run.font.color.rgb = RGBColor(128, 128, 128)
                            run.font.size = Pt(10)

                        # 添加中文释义
                        if meaning:
                            run = p.add_run(f" {meaning}")
                            run.font.size = Pt(11)

                        # 添加例句
                        example = word.get("example_sentence")
                        if example:
                            doc.add_paragraph(f"  例句: {example}", style='List Bullet 2')

        # 语法点
        if lesson_plan.grammar_points:
            doc.add_heading("三、语法要点", level=1)
            grammar_points = lesson_plan.grammar_points if isinstance(lesson_plan.grammar_points, list) else []

            for idx, gp in enumerate(grammar_points[:10], 1):
                doc.add_heading(f"{idx}. {gp.get('name', '语法点')}", level=2)

                if gp.get("description"):
                    doc.add_paragraph(gp["description"])

                if gp.get("rule"):
                    p = doc.add_paragraph()
                    run = p.add_run("规则: ")
                    run.bold = True
                    p.add_run(gp["rule"])

                if gp.get("examples"):
                    doc.add_paragraph("例句:", style='List Bullet')
                    for ex in gp["examples"][:3]:
                        doc.add_paragraph(f"  {ex}", style='List Bullet 2')

        # 教学流程
        if lesson_plan.teaching_structure:
            doc.add_heading("四、教学流程", level=1)
            structure = lesson_plan.teaching_structure

            phases = ["warm_up", "presentation", "practice", "production", "summary", "homework"]
            phase_names = {
                "warm_up": "热身环节",
                "presentation": "新知呈现",
                "practice": "练习环节",
                "production": "产出环节",
                "summary": "课堂总结",
                "homework": "课后作业"
            }

            for phase in phases:
                if phase in structure:
                    phase_data = structure[phase]
                    if isinstance(phase_data, dict) and phase_data:
                        doc.add_heading(phase_names.get(phase, phase), level=2)

                        if phase_data.get("duration"):
                            doc.add_paragraph(f"时长: {phase_data['duration']} 分钟")

                        if phase_data.get("description"):
                            doc.add_paragraph(phase_data["description"])

                        if phase_data.get("activities"):
                            doc.add_paragraph("活动:", style='List Bullet')
                            for activity in phase_data["activities"][:5]:
                                doc.add_paragraph(f"  {activity}", style='List Bullet 2')

        # 练习题
        if lesson_plan.exercises:
            doc.add_heading("五、练习题", level=1)
            exercises = lesson_plan.exercises

            for ex_type, ex_list in exercises.items():
                if ex_list:
                    type_name = {
                        "multiple_choice": "选择题",
                        "fill_blank": "填空题",
                        "matching": "匹配题",
                        "essay": "写作题"
                    }.get(ex_type, ex_type)
                    doc.add_heading(type_name, level=2)

                    for idx, ex in enumerate(ex_list[:10], 1):
                        doc.add_paragraph(f"{idx}. {ex.get('question', '')}")

                        options = ex.get("options")
                        if options:
                            for opt in options:
                                doc.add_paragraph(f"  {opt}", style='List Bullet')

                        # 添加答案
                        if ex.get("correct_answer"):
                            p = doc.add_paragraph()
                            run = p.add_run("答案: ")
                            run.bold = True
                            p.add_run(ex["correct_answer"])

                        if ex.get("explanation"):
                            doc.add_paragraph(f"解析: {ex['explanation']}", style='List Bullet')

        # PPT 大纲
        if lesson_plan.ppt_outline:
            doc.add_heading("六、PPT 大纲", level=1)
            ppt = lesson_plan.ppt_outline if isinstance(lesson_plan.ppt_outline, list) else []

            for slide in ppt[:20]:
                doc.add_heading(f"幻灯片 {slide.get('slide_number', '')}: {slide.get('title', '')}", level=2)

                content = slide.get("content", [])
                if content:
                    for item in content:
                        doc.add_paragraph(f"- {item}", style='List Bullet')

                if slide.get("notes"):
                    p = doc.add_paragraph()
                    run = p.add_run("备注: ")
                    run.font.color.rgb = RGBColor(128, 128, 128)
                    p.add_run(slide["notes"])

        # 教学反思
        if lesson_plan.teaching_notes:
            doc.add_heading("七、教学反思", level=1)
            doc.add_paragraph(lesson_plan.teaching_notes)

        # 保存到内存
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)

        filename = f"{lesson_plan.title}.docx"
        return file_stream.read(), filename

    @staticmethod
    async def export_to_ppt(lesson_plan: LessonPlan) -> tuple[bytes, str]:
        """
        导出教案为 PowerPoint 演示文稿

        Args:
            lesson_plan: 教案对象

        Returns:
            tuple[bytes, str]: (文件内容, 文件名)
        """
        prs = Presentation()

        # 设置幻灯片大小为 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # 封面页
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = lesson_plan.title
        subtitle.text = f"{lesson_plan.topic} | {lesson_plan.level} | {lesson_plan.duration}分钟"

        # 课程信息页
        info_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = info_slide.shapes.title
        title_shape.text = "课程信息"

        # 使用文本框添加信息
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)

        info_box = info_slide.shapes.add_textbox(left, top, width, height)
        text_frame = info_box.text_frame
        text_frame.word_wrap = True

        info_text = f"""教学主题: {lesson_plan.topic}
CEFR 等级: {lesson_plan.level}
课程时长: {lesson_plan.duration} 分钟
目标考试: {lesson_plan.target_exam or '通用'}"""

        p = text_frame.paragraphs[0]
        p.text = info_text
        p.font.size = Pt(18)

        # 教学目标页
        if lesson_plan.objectives:
            obj_slide = prs.slides.add_slide(prs.slide_layouts[1])
            obj_slide.shapes.title.text = "教学目标"

            content_box = obj_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
            text_frame = content_box.text_frame

            objectives = lesson_plan.objectives
            if objectives.get("language_knowledge"):
                p = text_frame.paragraphs[0]
                p.text = "知识目标"
                p.font.bold = True
                p.font.size = Pt(20)

                for obj in objectives["language_knowledge"][:5]:
                    p = text_frame.add_paragraph()
                    p.text = f"• {obj}"
                    p.font.size = Pt(16)
                    p.level = 1

        # 词汇页
        if lesson_plan.vocabulary:
            vocab_slide = prs.slides.add_slide(prs.slide_layouts[1])
            vocab_slide.shapes.title.text = "核心词汇"

            content_box = vocab_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
            text_frame = content_box.text_frame

            vocabulary = lesson_plan.vocabulary
            word_count = 0
            max_words_per_slide = 8

            for pos_type, words in vocabulary.items():
                if words and word_count < max_words_per_slide:
                    for word in words[:max_words_per_slide - word_count]:
                        p = text_frame.paragraphs[0] if word_count == 0 else text_frame.add_paragraph()
                        p.text = f"{word.get('word', '')} - {word.get('meaning_cn', '')}"
                        p.font.size = Pt(18)
                        word_count += 1

        # 语法点页
        if lesson_plan.grammar_points:
            for idx, gp in enumerate(list(lesson_plan.grammar_points)[:3]):
                gp_slide = prs.slides.add_slide(prs.slide_layouts[1])
                gp_slide.shapes.title.text = f"语法: {gp.get('name', '')}"

                content_box = gp_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.5))
                text_frame = content_box.text_frame

                if gp.get("description"):
                    p = text_frame.paragraphs[0]
                    p.text = gp["description"]
                    p.font.size = Pt(16)

                if gp.get("rule"):
                    p = text_frame.add_paragraph()
                    p.text = f"规则: {gp['rule']}"
                    p.font.size = Pt(14)
                    p.font.bold = True

                if gp.get("examples"):
                    p = text_frame.add_paragraph()
                    p.text = "例句:"
                    p.font.size = Pt(14)

                    for ex in gp["examples"][:2]:
                        p = text_frame.add_paragraph()
                        p.text = f"• {ex}"
                        p.font.size = Pt(12)
                        p.level = 1

        # 教学流程页
        if lesson_plan.teaching_structure:
            structure_slide = prs.slides.add_slide(prs.slide_layouts[1])
            structure_slide.shapes.title.text = "教学流程"

            content_box = structure_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
            text_frame = content_box.text_frame

            structure = lesson_plan.teaching_structure
            phases = ["warm_up", "presentation", "practice", "production", "summary"]

            phase_names = {
                "warm_up": "热身",
                "presentation": "新知呈现",
                "practice": "练习",
                "production": "产出",
                "summary": "总结"
            }

            for phase in phases:
                if phase in structure:
                    phase_data = structure[phase]
                    if isinstance(phase_data, dict):
                        p = text_frame.paragraphs[0] if phase == phases[0] else text_frame.add_paragraph()
                        duration = phase_data.get("duration", 0)
                        p.text = f"• {phase_names.get(phase, phase)} ({duration}分钟)"
                        p.font.size = Pt(16)

        # PPT 大纲页
        if lesson_plan.ppt_outline and isinstance(lesson_plan.ppt_outline, list):
            for slide_data in lesson_plan.ppt_outline[:10]:
                content_slide = prs.slides.add_slide(prs.slide_layouts[1])
                content_slide.shapes.title.text = slide_data.get("title", "")

                content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.5))
                text_frame = content_box.text_frame

                content = slide_data.get("content", [])
                for item in content:
                    p = text_frame.paragraphs[0] if item == content[0] else text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(16)

        # 保存到内存
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)

        filename = f"{lesson_plan.title}.pptx"
        return file_stream.read(), filename

    @staticmethod
    async def export_to_pdf(lesson_plan: LessonPlan) -> tuple[bytes, str]:
        """
        导出教案为 PDF 文档

        Args:
            lesson_plan: 教案对象

        Returns:
            tuple[bytes, str]: (文件内容, 文件名)
        """
        file_stream = io.BytesIO()
        doc = SimpleDocTemplate(file_stream, pagesize=A4)

        # 样式
        styles = getSampleStyleSheet()

        # 添加自定义样式
        styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=styles['Heading1'],
            fontSize=20,
            textColor=colors.HexColor('#2C3E50'),
            spaceAfter=30,
            alignment=1  # 居中
        ))

        styles.add(ParagraphStyle(
            name='SectionHeader',
            parent=styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#34495E'),
            spaceAfter=12,
            spaceBefore=12
        ))

        # 构建内容
        story = []

        # 标题
        story.append(Paragraph(lesson_plan.title, styles['CustomTitle']))
        story.append(Spacer(1, 0.2 * inch))

        # 基本信息表格
        info_data = [
            ["教学主题", lesson_plan.topic],
            ["CEFR 等级", lesson_plan.level],
            ["课程时长", f"{lesson_plan.duration} 分钟"],
            ["目标考试", lesson_plan.target_exam or "通用"],
        ]

        info_table = Table(info_data, colWidths=[1.5 * inch, 4 * inch])
        info_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.grey),
            ('TEXTCOLOR', (0, 0), (0, -1), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ]))
        story.append(info_table)
        story.append(Spacer(1, 0.3 * inch))

        # 教学目标
        if lesson_plan.objectives:
            story.append(Paragraph("教学目标", styles['SectionHeader']))
            objectives = lesson_plan.objectives

            if objectives.get("language_knowledge"):
                story.append(Paragraph("<b>知识目标</b>", styles['Normal']))
                for obj in objectives["language_knowledge"]:
                    story.append(Paragraph(f"• {obj}", styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))

        # 核心词汇
        if lesson_plan.vocabulary:
            story.append(PageBreak())
            story.append(Paragraph("核心词汇", styles['SectionHeader']))

            vocabulary = lesson_plan.vocabulary
            for pos_type, words in vocabulary.items():
                if words:
                    pos_name = {"noun": "名词", "verb": "动词"}.get(pos_type, pos_type)
                    story.append(Paragraph(f"<b>{pos_name}</b>", styles['Normal']))

                    for word in words[:15]:
                        word_text = word.get("word", "")
                        meaning = word.get("meaning_cn", "")
                        story.append(Paragraph(f"<b>{word_text}</b> - {meaning}", styles['Normal']))

                    story.append(Spacer(1, 0.1 * inch))

        # 语法点
        if lesson_plan.grammar_points:
            story.append(PageBreak())
            story.append(Paragraph("语法要点", styles['SectionHeader']))

            grammar_points = lesson_plan.grammar_points if isinstance(lesson_plan.grammar_points, list) else []
            for gp in grammar_points[:5]:
                story.append(Paragraph(f"<b>{gp.get('name', '')}</b>", styles['Normal']))
                if gp.get("description"):
                    story.append(Paragraph(gp["description"], styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))

        # 教学流程
        if lesson_plan.teaching_structure:
            story.append(PageBreak())
            story.append(Paragraph("教学流程", styles['SectionHeader']))

            structure = lesson_plan.teaching_structure
            phases = ["warm_up", "presentation", "practice", "production", "summary"]

            phase_names = {
                "warm_up": "热身",
                "presentation": "新知呈现",
                "practice": "练习",
                "production": "产出",
                "summary": "总结"
            }

            for phase in phases:
                if phase in structure:
                    phase_data = structure[phase]
                    if isinstance(phase_data, dict):
                        duration = phase_data.get("duration", 0)
                        story.append(Paragraph(f"<b>{phase_names.get(phase, phase)}</b> ({duration}分钟)", styles['Normal']))

                        if phase_data.get("description"):
                            story.append(Paragraph(phase_data["description"], styles['Normal']))
                        story.append(Spacer(1, 0.1 * inch))

        # 练习题
        if lesson_plan.exercises:
            story.append(PageBreak())
            story.append(Paragraph("练习题", styles['SectionHeader']))

            exercises = lesson_plan.exercises
            for ex_type, ex_list in exercises.items():
                if ex_list and isinstance(ex_list, list):
                    type_name = {"multiple_choice": "选择题"}.get(ex_type, ex_type)
                    story.append(Paragraph(f"<b>{type_name}</b>", styles['Normal']))

                    for idx, ex in enumerate(ex_list[:5], 1):
                        story.append(Paragraph(f"{idx}. {ex.get('question', '')}", styles['Normal']))

                        if ex.get("correct_answer"):
                            story.append(Paragraph(f"<i>答案: {ex['correct_answer']}</i>", styles['Normal']))

                        story.append(Spacer(1, 0.1 * inch))

        # 构建 PDF
        doc.build(story)
        file_stream.seek(0)

        filename = f"{lesson_plan.title}.pdf"
        return file_stream.read(), filename


# 创建全局单例
_export_service: Optional[ExportService] = None


def get_export_service() -> ExportService:
    """获取导出服务单例"""
    global _export_service
    if _export_service is None:
        _export_service = ExportService()
    return _export_service
